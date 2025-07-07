"""
多格式文档处理器
支持各种文档格式的内容提取、预览和分析
"""

import os
import io
import csv
import json
import yaml
import xml.etree.ElementTree as ET
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
import logging

# 文档处理库
try:
    from docx import Document as DocxDocument
    from openpyxl import load_workbook
    from pptx import Presentation
    import PyPDF2
    import pdfplumber
    from PIL import Image
    # import easyocr  # OCR库，需要安装: pip install easyocr
except ImportError as e:
    logging.warning(f"某些文档处理库未安装: {e}")

from app.models.knowledge import DocumentType

logger = logging.getLogger(__name__)


class DocumentProcessorError(Exception):
    """文档处理器异常"""
    pass


class DocumentProcessor:
    """多格式文档处理器"""
    
    def __init__(self):
        """初始化文档处理器"""
        self.supported_types = {
            # 文档类型
            DocumentType.PDF: self._process_pdf,
            DocumentType.DOCX: self._process_docx,
            DocumentType.DOC: self._process_doc,
            DocumentType.TXT: self._process_txt,
            DocumentType.MD: self._process_txt,
            DocumentType.RTF: self._process_rtf,
            
            # 表格类型
            DocumentType.XLSX: self._process_xlsx,
            DocumentType.XLS: self._process_xls,
            DocumentType.CSV: self._process_csv,
            
            # 演示文稿类型
            DocumentType.PPTX: self._process_pptx,
            DocumentType.PPT: self._process_ppt,
            
            # 数据格式
            DocumentType.JSON: self._process_json,
            DocumentType.XML: self._process_xml,
            DocumentType.YAML: self._process_yaml,
            DocumentType.YML: self._process_yaml,
            
            # 图片类型 (OCR)
            DocumentType.PNG: self._process_image_ocr,
            DocumentType.JPG: self._process_image_ocr,
            DocumentType.JPEG: self._process_image_ocr,
            DocumentType.GIF: self._process_image_ocr,
            DocumentType.WEBP: self._process_image_ocr,
            
            # 网页格式
            DocumentType.HTML: self._process_html,
            DocumentType.HTM: self._process_html,
        }
        
        # 初始化OCR引擎（延迟加载）
        self._ocr_reader = None
    
    def is_supported(self, document_type: DocumentType) -> bool:
        """检查是否支持该文档类型"""
        return document_type in self.supported_types
    
    def get_supported_types(self) -> List[DocumentType]:
        """获取所有支持的文档类型"""
        return list(self.supported_types.keys())
    
    def process_document(self, file_path: str, document_type: DocumentType) -> Dict[str, Any]:
        """
        处理文档，提取内容和元数据
        
        Args:
            file_path: 文档文件路径
            document_type: 文档类型
            
        Returns:
            包含文档内容和元数据的字典
        """
        if not self.is_supported(document_type):
            raise DocumentProcessorError(f"不支持的文档类型: {document_type}")
        
        if not os.path.exists(file_path):
            raise DocumentProcessorError(f"文件不存在: {file_path}")
        
        try:
            processor_func = self.supported_types[document_type]
            return processor_func(file_path)
        except Exception as e:
            logger.error(f"处理文档失败 {file_path}: {e}")
            raise DocumentProcessorError(f"处理文档失败: {e}")
    
    def get_document_preview(self, file_path: str, document_type: DocumentType, max_length: int = 500) -> str:
        """
        获取文档预览内容
        
        Args:
            file_path: 文档文件路径
            document_type: 文档类型
            max_length: 预览最大字符数
            
        Returns:
            文档预览文本
        """
        try:
            result = self.process_document(file_path, document_type)
            content = result.get('content', '')
            
            if len(content) <= max_length:
                return content
            
            return content[:max_length] + "..."
        except Exception as e:
            logger.error(f"获取文档预览失败: {e}")
            return f"预览生成失败: {e}"
    
    # ========== PDF 处理 ==========
    
    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """处理PDF文档"""
        try:
            content = ""
            metadata = {}
            
            # 使用pdfplumber提取文本（更好的文本提取效果）
            with pdfplumber.open(file_path) as pdf:
                metadata.update({
                    'page_count': len(pdf.pages),
                    'title': pdf.metadata.get('Title', ''),
                    'author': pdf.metadata.get('Author', ''),
                    'creator': pdf.metadata.get('Creator', ''),
                    'subject': pdf.metadata.get('Subject', ''),
                })
                
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        content += page_text + "\n"
            
            return {
                'content': content.strip(),
                'metadata': metadata,
                'processor': 'pdfplumber'
            }
        except Exception as e:
            # 降级到PyPDF2
            return self._process_pdf_fallback(file_path)
    
    def _process_pdf_fallback(self, file_path: str) -> Dict[str, Any]:
        """PDF处理降级方案"""
        content = ""
        metadata = {}
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            metadata.update({
                'page_count': len(pdf_reader.pages),
            })
            
            if pdf_reader.metadata:
                metadata.update({
                    'title': pdf_reader.metadata.get('/Title', ''),
                    'author': pdf_reader.metadata.get('/Author', ''),
                    'creator': pdf_reader.metadata.get('/Creator', ''),
                    'subject': pdf_reader.metadata.get('/Subject', ''),
                })
            
            for page in pdf_reader.pages:
                content += page.extract_text() + "\n"
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'processor': 'PyPDF2'
        }
    
    # ========== Word 文档处理 ==========
    
    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """处理DOCX文档"""
        doc = DocxDocument(file_path)
        
        # 提取文本内容
        content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content.append(paragraph.text)
        
        # 提取表格内容
        tables_content = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            tables_content.append(table_data)
        
        # 提取元数据
        metadata = {
            'paragraph_count': len(doc.paragraphs),
            'table_count': len(doc.tables),
        }
        
        # 尝试获取文档属性
        try:
            core_props = doc.core_properties
            metadata.update({
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'created': str(core_props.created) if core_props.created else '',
                'modified': str(core_props.modified) if core_props.modified else '',
            })
        except:
            pass
        
        return {
            'content': '\n'.join(content),
            'tables': tables_content,
            'metadata': metadata,
            'processor': 'python-docx'
        }
    
    def _process_doc(self, file_path: str) -> Dict[str, Any]:
        """处理DOC文档（旧版Word格式）"""
        # DOC格式处理较复杂，这里返回基本信息
        return {
            'content': f"DOC格式文档: {os.path.basename(file_path)}\n需要转换为DOCX格式以获得更好的处理效果。",
            'metadata': {'format': 'DOC (Legacy)', 'requires_conversion': True},
            'processor': 'basic'
        }
    
    # ========== Excel 处理 ==========
    
    def _process_xlsx(self, file_path: str) -> Dict[str, Any]:
        """处理XLSX文档"""
        workbook = load_workbook(file_path, read_only=True, data_only=True)
        
        content = []
        sheets_data = {}
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_content = []
            
            # 读取工作表数据
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):  # 跳过空行
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    sheet_content.append(row_data)
            
            sheets_data[sheet_name] = sheet_content
            
            # 生成文本内容
            if sheet_content:
                content.append(f"工作表: {sheet_name}")
                for row in sheet_content[:10]:  # 只取前10行作为预览
                    content.append('\t'.join(row))
                content.append("")
        
        metadata = {
            'sheet_count': len(workbook.sheetnames),
            'sheet_names': workbook.sheetnames,
        }
        
        return {
            'content': '\n'.join(content),
            'sheets': sheets_data,
            'metadata': metadata,
            'processor': 'openpyxl'
        }
    
    def _process_xls(self, file_path: str) -> Dict[str, Any]:
        """处理XLS文档（旧版Excel格式）"""
        # XLS格式处理较复杂，这里返回基本信息
        return {
            'content': f"XLS格式文档: {os.path.basename(file_path)}\n建议转换为XLSX格式以获得更好的处理效果。",
            'metadata': {'format': 'XLS (Legacy)', 'requires_conversion': True},
            'processor': 'basic'
        }
    
    def _process_csv(self, file_path: str) -> Dict[str, Any]:
        """处理CSV文档"""
        content = []
        csv_data = []
        
        with open(file_path, 'r', encoding='utf-8-sig', newline='') as file:
            try:
                # 自动检测分隔符
                sample = file.read(1024)
                file.seek(0)
                sniffer = csv.Sniffer()
                delimiter = sniffer.sniff(sample).delimiter
            except:
                delimiter = ','
            
            reader = csv.reader(file, delimiter=delimiter)
            for i, row in enumerate(reader):
                csv_data.append(row)
                if i < 20:  # 只取前20行作为内容预览
                    content.append('\t'.join(row))
        
        metadata = {
            'row_count': len(csv_data),
            'column_count': len(csv_data[0]) if csv_data else 0,
            'delimiter': delimiter,
        }
        
        return {
            'content': '\n'.join(content),
            'data': csv_data,
            'metadata': metadata,
            'processor': 'csv'
        }
    
    # ========== PowerPoint 处理 ==========
    
    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """处理PPTX文档"""
        prs = Presentation(file_path)
        
        content = []
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            
            # 提取幻灯片中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            slides_content.append(slide_text)
            
            if slide_text:
                content.append(f"幻灯片 {i + 1}:")
                content.extend(slide_text)
                content.append("")
        
        metadata = {
            'slide_count': len(prs.slides),
        }
        
        return {
            'content': '\n'.join(content),
            'slides': slides_content,
            'metadata': metadata,
            'processor': 'python-pptx'
        }
    
    def _process_ppt(self, file_path: str) -> Dict[str, Any]:
        """处理PPT文档（旧版PowerPoint格式）"""
        return {
            'content': f"PPT格式文档: {os.path.basename(file_path)}\n建议转换为PPTX格式以获得更好的处理效果。",
            'metadata': {'format': 'PPT (Legacy)', 'requires_conversion': True},
            'processor': 'basic'
        }
    
    # ========== 纯文本处理 ==========
    
    def _process_txt(self, file_path: str) -> Dict[str, Any]:
        """处理TXT和MD文档"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        except UnicodeDecodeError:
            # 如果UTF-8解码失败，尝试其他编码
            for encoding in ['gbk', 'gb2312', 'latin-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    break
                except:
                    continue
            else:
                content = "无法解码文件内容"
        
        lines = content.split('\n')
        metadata = {
            'line_count': len(lines),
            'char_count': len(content),
            'word_count': len(content.split()) if content else 0,
        }
        
        return {
            'content': content,
            'metadata': metadata,
            'processor': 'text'
        }
    
    def _process_rtf(self, file_path: str) -> Dict[str, Any]:
        """处理RTF文档"""
        # RTF处理需要专门的库，这里简化处理
        return {
            'content': f"RTF格式文档: {os.path.basename(file_path)}\n需要专门的RTF处理库来提取内容。",
            'metadata': {'format': 'RTF', 'requires_special_library': True},
            'processor': 'basic'
        }
    
    # ========== 数据格式处理 ==========
    
    def _process_json(self, file_path: str) -> Dict[str, Any]:
        """处理JSON文档"""
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
        
        # 将JSON转换为可读文本
        content = json.dumps(data, indent=2, ensure_ascii=False)
        
        metadata = {
            'type': type(data).__name__,
            'size': len(str(data)),
        }
        
        if isinstance(data, dict):
            metadata['keys'] = list(data.keys())[:10]  # 只显示前10个键
        elif isinstance(data, list):
            metadata['length'] = len(data)
        
        return {
            'content': content,
            'data': data,
            'metadata': metadata,
            'processor': 'json'
        }
    
    def _process_xml(self, file_path: str) -> Dict[str, Any]:
        """处理XML文档"""
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            # 提取XML内容为文本
            content = ET.tostring(root, encoding='unicode')
            
            metadata = {
                'root_tag': root.tag,
                'element_count': len(list(root.iter())),
            }
            
            return {
                'content': content,
                'metadata': metadata,
                'processor': 'xml'
            }
        except ET.ParseError as e:
            return {
                'content': f"XML解析错误: {e}",
                'metadata': {'parse_error': str(e)},
                'processor': 'xml'
            }
    
    def _process_yaml(self, file_path: str) -> Dict[str, Any]:
        """处理YAML文档"""
        with open(file_path, 'r', encoding='utf-8') as file:
            data = yaml.safe_load(file)
        
        # 将YAML转换为可读文本
        content = yaml.dump(data, default_flow_style=False, allow_unicode=True)
        
        metadata = {
            'type': type(data).__name__,
        }
        
        if isinstance(data, dict):
            metadata['keys'] = list(data.keys())[:10]
        elif isinstance(data, list):
            metadata['length'] = len(data)
        
        return {
            'content': content,
            'data': data,
            'metadata': metadata,
            'processor': 'yaml'
        }
    
    # ========== 图片OCR处理 ==========
    
    def _process_image_ocr(self, file_path: str) -> Dict[str, Any]:
        """处理图片文档（OCR文字识别）"""
        try:
            # 延迟初始化OCR
            # if self._ocr_reader is None:
            #     self._ocr_reader = easyocr.Reader(['ch_sim', 'en'])  # 中英文识别
            
            # 进行OCR识别
            # results = self._ocr_reader.readtext(file_path)
            
            # 提取文本内容
            # content = []
            # for (bbox, text, confidence) in results:
            #     if confidence > 0.5:  # 只保留置信度较高的文本
            #         content.append(text)
            
            # 获取图片基本信息
            with Image.open(file_path) as img:
                width, height = img.size
                mode = img.mode
            
            metadata = {
                'width': width,
                'height': height,
                'mode': mode,
                'ocr_available': False,  # 暂时关闭OCR功能
            }
            
            return {
                'content': f"图片文档: {os.path.basename(file_path)}\n尺寸: {width}x{height}\n模式: {mode}\n\nOCR功能暂未启用，需要安装easyocr库",
                'metadata': metadata,
                'processor': 'pillow'
            }
        except Exception as e:
            # OCR失败时的降级处理
            return {
                'content': f"图片文档: {os.path.basename(file_path)}\n处理失败: {e}",
                'metadata': {'process_error': str(e)},
                'processor': 'basic'
            }
    
    # ========== HTML处理 ==========
    
    def _process_html(self, file_path: str) -> Dict[str, Any]:
        """处理HTML文档"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 提取文本内容
            text_content = soup.get_text(separator='\n', strip=True)
            
            # 提取元数据
            metadata = {
                'title': soup.title.string if soup.title else '',
                'links_count': len(soup.find_all('a')),
                'images_count': len(soup.find_all('img')),
                'tables_count': len(soup.find_all('table')),
            }
            
            # 提取meta标签信息
            meta_tags = soup.find_all('meta')
            for meta in meta_tags:
                if meta.get('name'):
                    metadata[f"meta_{meta.get('name')}"] = meta.get('content', '')
            
            return {
                'content': text_content,
                'html': html_content,
                'metadata': metadata,
                'processor': 'beautifulsoup'
            }
        except ImportError:
            # 如果没有安装BeautifulSoup，使用基本处理
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            return {
                'content': content,
                'metadata': {'requires_beautifulsoup': True},
                'processor': 'basic'
            }


# 创建全局文档处理器实例
document_processor = DocumentProcessor()


def get_document_processor() -> DocumentProcessor:
    """获取文档处理器实例"""
    return document_processor 